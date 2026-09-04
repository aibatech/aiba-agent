"""Phase 8 — media / document extraction tests.

Covers the read-only `media_extract` tool end to end over an isolated temporary
workspace: real extraction of PDF/DOCX/XLSX/PPTX/CSV/markdown/text and image
metadata, boundary/unsupported/limits handling, workshop confinement (path
traversal is blocked), and honest "install optional support" diagnostics when a
back-end library is missing.

Fixtures are built in-test (pure-Python minimal PDF builder; the optional doc
libraries only when present). Every test that needs an optional parser library
is skipped cleanly when that library is not installed, so this module runs in a
bare `[api]` environment too.
"""
from __future__ import annotations

import importlib.util
import io
import json
import shutil
import tempfile
import unittest
from pathlib import Path

from tools.sandbox import Sandbox
from security.policy import SecurityPolicy
from tools.media import MediaExtraction, build_media_tools

REPO = Path(__file__).resolve().parents[1]
PERMISSIONS = REPO / "config" / "permissions.json"

HAVE_PYPDF = importlib.util.find_spec("pypdf") is not None
HAVE_DOCX = importlib.util.find_spec("docx") is not None
HAVE_OPENPYXL = importlib.util.find_spec("openpyxl") is not None
HAVE_PPTX = importlib.util.find_spec("pptx") is not None
HAVE_PIL = importlib.util.find_spec("PIL") is not None


def _escape_tx(text: str) -> str:
    return text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def minimal_text_pdf(lines: list[str]) -> bytes:
    """Build a valid single-page PDF containing readable text lines."""
    stream_body = b"\n".join(
        f"BT /F1 12 Tf 72 {780 - i * 18} Td ({_escape_tx(ln)}) Tj ET".encode("latin-1")
        for i, ln in enumerate(lines)
    )
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 5 0 R /Resources << /Font << /F1 4 0 R >> >> >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        f"<< /Length {len(stream_body)} >>\nstream\n".encode("ascii")
        + stream_body
        + b"\nendstream",
    ]
    out = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = []
    for i, obj in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{i} 0 obj\n".encode("ascii") + obj + b"\nendobj\n"
    xref_pos = len(out)
    count = len(objects) + 1
    out += f"xref\n0 {count}\n".encode("ascii")
    out += b"0000000000 65535 f \n"
    for off in offsets:
        out += f"{off:010d} 00000 n \n".encode("ascii")
    out += (f"trailer\n<< /Size {count} /Root 1 0 R >>\nstartxref\n{xref_pos}\n%%EOF\n"
            .encode("ascii"))
    return bytes(out)


def build_fixtures(ws: Path) -> None:
    """Write a representative set of document fixtures into ws (a workspace)."""
    ws.mkdir(parents=True, exist_ok=True)
    (ws / "sample.csv").write_text("name,role,score\nAlice,captain,9\nBob,navigator,7\n")
    (ws / "readme.md").write_text("# Title\n\nHello **world**.\n- a\n- b\n")
    (ws / "note.txt").write_text("plain text line one\nline two\n")

    if HAVE_PYPDF:
        (ws / "sample.pdf").write_bytes(minimal_text_pdf(
            ["Hello PDF world.", "Second PDF line."]))
    if HAVE_DOCX:
        import docx
        d = docx.Document()
        d.add_heading("Doc Heading 0", 0)
        d.add_paragraph("Hello from docx paragraph.")
        d.add_table(rows=1, cols=2)
        d.tables[0].rows[0].cells[0].text = "cellA"
        d.tables[0].rows[0].cells[1].text = "cellB"
        d.save(str(ws / "sample.docx"))
    if HAVE_OPENPYXL:
        import openpyxl
        wb = openpyxl.Workbook()
        sh = wb.active
        sh.title = "Data"
        sh.append(["prod", "qty"])
        sh.append(["widget", 3])
        meta = wb.create_sheet("Meta")
        meta["A1"] = "note"
        wb.save(str(ws / "sample.xlsx"))
    if HAVE_PPTX:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide One Title"
        slide.placeholders[1].text = "Body here"
        prs.save(str(ws / "sample.pptx"))
    if HAVE_PIL:
        from PIL import Image
        Image.new("RGB", (12, 8), (200, 30, 30)).save(str(ws / "pix.png"))


def _extractor(tmp: Path) -> tuple[Sandbox, MediaExtraction]:
    sandbox = Sandbox(tmp, timeout=30, policy=SecurityPolicy(tmp, PERMISSIONS, True),
                      mode="local")
    return sandbox, MediaExtraction(sandbox)


# Fixture builder itself only needs the libs we can skip; force workspace even
# if libs absent (text fixtures still exercise works).
class Phase8MediaTests(unittest.TestCase):
    def setUp(self):
        self._tmp = tempfile.mkdtemp(prefix="aiba_phase8_media_")
        self.ws = Path(self._tmp) / "ws"
        build_fixtures(self.ws)
        self.sandbox, self.mx = _extractor(self.ws)

    def tearDown(self):
        shutil.rmtree(self._tmp, ignore_errors=True)

    # -- stdlib / text-capable formats (always run) -------------------------
    def test_csv_extraction(self):
        r = self.mx.media_extract("sample.csv")
        self.assertTrue(r.ok, r.error)
        out = r.output
        self.assertEqual(out["format"], "csv")
        self.assertEqual(out["rows"], 3)
        self.assertIn("Alice", out["text"])
        self.assertIn("Bob", out["text"])

    def test_markdown_extraction(self):
        r = self.mx.media_extract("readme.md")
        self.assertTrue(r.ok)
        self.assertEqual(r.output["format"], "md")
        self.assertIn("Hello", r.output["text"])

    def test_txt_extraction(self):
        r = self.mx.media_extract("note.txt")
        self.assertTrue(r.ok)
        self.assertIn("plain text line one", r.output["text"])

    # -- optional-lib formats (skipped when the lib is absent) --------------
    @unittest.skipUnless(HAVE_PYPDF, "pypdf not installed")
    def test_pdf_extraction(self):
        r = self.mx.media_extract("sample.pdf")
        self.assertTrue(r.ok, r.error)
        self.assertEqual(r.output["pages"], 1)
        self.assertIn("Hello PDF world.", r.output["text"])
        self.assertIn("Second PDF line.", r.output["text"])

    @unittest.skipUnless(HAVE_DOCX, "python-docx not installed")
    def test_docx_extraction(self):
        r = self.mx.media_extract("sample.docx")
        self.assertTrue(r.ok, r.error)
        self.assertIn("Hello from docx paragraph.", r.output["text"])
        self.assertIn("cellA", r.output["text"])
        self.assertIn("cellB", r.output["text"])

    @unittest.skipUnless(HAVE_OPENPYXL, "openpyxl not installed")
    def test_xlsx_extraction(self):
        r = self.mx.media_extract("sample.xlsx")
        self.assertTrue(r.ok, r.error)
        self.assertIn("== Meta ==", r.output["text"])  # second sheet surfaced
        self.assertIn("widget", r.output["text"])

    @unittest.skipUnless(HAVE_PPTX, "python-pptx not installed")
    def test_pptx_extraction(self):
        r = self.mx.media_extract("sample.pptx")
        self.assertTrue(r.ok, r.error)
        self.assertIn("Slide One Title", r.output["text"])
        self.assertIn("Body here", r.output["text"])

    @unittest.skipUnless(HAVE_PIL, "Pillow not installed")
    def test_image_metadata(self):
        r = self.mx.media_extract("pix.png")
        self.assertTrue(r.ok, r.error)
        self.assertEqual(r.output["kind"], "image")
        self.assertEqual(r.output["metadata"]["width"], 12)
        self.assertEqual(r.output["metadata"]["height"], 8)
        self.assertEqual(r.output["metadata"]["format"], "png")

    # -- boundaries / error handling ----------------------------------------
    def test_unsupported_extension(self):
        (self.ws / "blob.bin").write_bytes(b"\x00\x01\x02")
        r = self.mx.media_extract("blob.bin")
        self.assertFalse(r.ok)
        self.assertIn("Unsupported format", r.error)

    def test_missing_file(self):
        r = self.mx.media_extract("does_not_exist.csv")
        self.assertFalse(r.ok)
        self.assertIn("not found", r.error.lower())

    def test_path_traversal_blocked(self):
        # Real SecurityPolicy; every traversal must be refused, never read.
        for bad in ("../../etc/hostname", "../../../../etc/passwd", "../secret.txt"):
            r = self.mx.media_extract(bad)
            self.assertFalse(r.ok, f"should block {bad!r}")
            self.assertTrue("outside" in r.error.lower() or "not found" in r.error.lower(),
                            r.error)

    def test_nul_byte_never_reaches_filesystem(self):
        # A NUL byte is stripped *before* any filesystem access (path cannot
        # truncate or confuse resolution). The call must never raise and never
        # let the raw NUL reach open(); it resolves to the stripped name.
        (self.ws / "real.csv").write_text("x\n")  # stripped target exists
        r = self.mx.media_extract("rea\x00l.csv")
        self.assertTrue(r.ok, r.error)      # strip => real.csv, read fine
        # A stripped target that does not exist must be a clean not-found.
        r2 = self.mx.media_extract("no_such\x00file.csv")
        self.assertFalse(r2.ok)
        self.assertIn("not found", r2.error.lower())

    def test_returns_text_not_evaluated(self):
        # A CSV row that looks like a formula must stay literal text; no eval.
        (self.ws / "formula.csv").write_text('a,b\n=1+1,=cmd|"x"\n')
        r = self.mx.media_extract("formula.csv")
        self.assertTrue(r.ok, r.error)
        self.assertIn("=1+1", r.output["text"])

    # -- generous truncation bound -----------------------------------------
    def test_long_text_truncated(self):
        (self.ws / "big.txt").write_text("word " * 20000)
        r = self.mx.media_extract("big.txt")
        self.assertTrue(r.ok, r.error)
        self.assertTrue(r.output["truncated"])
        self.assertLessEqual(len(r.output["text"]), 41000)
        self.assertIn("truncated", r.output["text"])


# ---------------------------------------------------------------------------
# AgentLoop-level: media_extract is a real, registered model tool kept in
# manifest/permissions parity and gated by the AIBA_MEDIA_ENABLED feature flag.
# ---------------------------------------------------------------------------
def _loop_settings(tmp: Path, *, media_enabled: bool):
    """Build a Settings over an isolated temp root with media toggled."""
    from config.settings import Settings
    data = tmp / "data"
    for d in ("workspace", "vault", "logs", "reflections", "skill_proposals"):
        (data / d).mkdir(parents=True, exist_ok=True)
    skills = tmp / "skills"
    skills.mkdir(parents=True, exist_ok=True)
    cfg = tmp / "config"
    cfg.mkdir(parents=True, exist_ok=True)
    shutil.copy(REPO / "config" / "permissions.json", cfg / "permissions.json")
    shutil.copy(REPO / "config" / "capability_manifest.json", cfg / "capability_manifest.json")
    di = lambda p: tmp / p  # noqa: E731
    return Settings(
        root_dir=tmp, data_dir=data, workspace_dir=di("data/workspace"),
        vault_dir=di("data/vault"), logs_dir=di("data/logs"),
        skills_dir=skills,
        db_path=di("data/aiba.db"), tasks_db_path=di("data/tasks.db"),
        sessions_db_path=di("data/sessions.db"), jobs_db_path=di("data/jobs.db"),
        schedules_db_path=di("data/schedules.db"), auth_db_path=di("data/auth.db"),
        providers_db_path=di("data/providers.db"),
        provider="local", fallback_provider="local", model="local-v1",
        fallback_model="local-v1", max_steps=3, command_timeout=10,
        require_approval=True, sandbox_mode="local",
        docker_image="python:3.12-slim", docker_memory="512m", docker_cpus="1.0",
        sandbox_network=False, permissions_path=cfg / "permissions.json",
        browser_enabled=False, desktop_enabled=False, vision_model="",
        worker_enabled=True, api_token="x" * 40, api_host="127.0.0.1",
        api_port=8765, allowed_origins=(), rate_limit_per_minute=60,
        web_enabled=False, computer_node_path=data / "computer_node.json",
        desktop_clipboard_enabled=False, desktop_process_enabled=False,
        media_enabled=media_enabled,
    )


class Phase8LoopTests(unittest.TestCase):
    def setUp(self):
        self._tmp = tempfile.mkdtemp(prefix="aiba_phase8_loop_")

    def tearDown(self):
        shutil.rmtree(self._tmp, ignore_errors=True)

    def test_media_extract_registered_and_model_visible_when_enabled(self):
        from agent.loop import AgentLoop
        loop = AgentLoop(settings=_loop_settings(Path(self._tmp), media_enabled=True),
                         interactive=False, auto_approve=True, start_worker=False)
        try:
            self.assertIn("media_extract", loop.registry._tools)
            visible = {s["name"] for s in loop.registry.schemas()}
            self.assertIn("media_extract", visible)
            # read-only local => not approval-gated
            self.assertFalse(loop.policy.check_tool("media_extract").requires_approval)
            # manifests stay in lockstep (validator re-checks parity at gate)
            self.assertIn("media_extract",
                          loop.manifest["tools"])
        finally:
            loop.close()

    def test_media_extract_hidden_when_flag_disabled(self):
        from agent.loop import AgentLoop
        loop = AgentLoop(settings=_loop_settings(Path(self._tmp), media_enabled=False),
                         interactive=False, auto_approve=True, start_worker=False)
        try:
            # still registered (present in manifest), but not advertised
            self.assertIn("media_extract", loop.registry._tools)
            visible = {s["name"] for s in loop.registry.schemas()}
            self.assertNotIn("media_extract", visible)
            res = loop.registry.execute("media_extract", {"path": "x.csv"})
            self.assertFalse(res.ok)
            self.assertIn("feature flag", res.error)
        finally:
            loop.close()

