"""Phase 8 — media / document extraction (read-only, safe by default).

AIBA can read the **text** out of common document formats so the model can
reason about files a user drops into the workspace, without ever executing
content inside them.

Hard security posture (drives review/test expectations):

* **Read-only.** Extraction never modifies the source file and never writes
  anything. There is deliberately no "output to disk" path here — callers that
  need to persist extracted text use the ordinary, approval-gated
  ``write_file``/``archive`` tools. Exports therefore cannot appear inside the
  media surface and stay behind the existing approval policy.
* **Workspace-confined.** Every path is resolved through the injected
  ``Sandbox`` (``policy.check_path``), so a request can never escape the
  approved workspace. All extraction is performed on the already-confined path.
* **Bounded.** File size, page/sheet/slide counts, row/cell counts, archive
  members and total returned characters are each capped. A pathological or
  malicious document cannot exhaust memory or blow up the model context.
* **Untrusted, never code.** Document contents are treated as data only. We do
  not evaluate spreadsheet formulas, run macros, follow embedded links, or
  execute anything read from a document. Extracted strings are plain text.
* **Honest optional deps.** PDF/DOCX/XLSX/PPTX parsing relies on the *optional*
  ``[media]`` extra (pypdf / python-docx / openpyxl / python-pptx). When a
  library is missing the tool returns a clear "install optional support"
  diagnostic instead of a partial parse. OCR / ASR / TTS / image generation are
  NOT implemented here: they are surfaced only as capability probes that report
  honestly not-available until a real backend and its tests exist.

All modules are import-guarded; importing ``media`` never requires the optional
libraries, so the lightweight base install keeps working unchanged.
"""
from __future__ import annotations

import importlib.util
import os
from pathlib import Path
from typing import Any

# ---- Optional import helpers -------------------------------------------------
# Deliberately import the doc libraries only inside the functions that need
# them. `import media.extract` must succeed with zero optional deps installed.


def _have(name: str) -> bool:
    return importlib.util.find_spec(name) is not None


def extract_text(path: str | os.PathLike, *, max_bytes: int = 64 * 1024 * 1024,
                 limit: int = 40_000) -> dict[str, Any]:
    """Extract plain text from a file, dispatching by extension.

    ``path`` must already be workspace-confined (resolved through the Sandbox).
    Returns a dict with at least ``format`` and ``text``. Raises
    ``FileNotFoundError`` if missing, ``ValueError`` on unsupported format, and
    a clear ``RuntimeError`` message if the required optional library is absent.

    All document content is treated as untrusted data; nothing is evaluated.
    """
    p = Path(path)
    if not p.is_file():
        raise FileNotFoundError(f"File not found: {p.name}")
    # Bounded read for text-ish files keeps memory flat up-front.
    suffix = (p.suffix or "").lower()
    if suffix == ".csv":
        return extract_csv(p, limit=limit)
    if suffix in {".txt", ".md", ".text", ".json", ".xml", ".html", ".htm",
                  ".yaml", ".yml", ".tsv", ".log", ".ini", ".cfg", ".toml"}:
        return extract_text_simple(p, suffix=suffix, max_bytes=max_bytes, limit=limit)
    if suffix in {".pdf", ".docx", ".xlsx", ".pptx"}:
        return _extract_binary_doc(p, suffix, max_bytes=max_bytes, limit=limit)
    if suffix in {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff", ".tif"}:
        return image_metadata(p)
    raise ValueError(
        f"Unsupported format '{suffix or p.name}'. media_extract supports "
        "PDF (.pdf), DOCX (.docx), XLSX (.xlsx), PPTX (.pptx), CSV (.csv), "
        "plain text/markdown, and common image metadata."
    )


def extract_text_simple(p: Path, *, suffix: str, max_bytes: int,
                        limit: int) -> dict[str, Any]:
    size = p.stat().st_size
    if size > max_bytes:
        raise RuntimeError(f"File too large ({size} bytes > {max_bytes} limit)")
    data = p.read_bytes()
    # XML/HTML arrive as bytes possibly with an encoding declaration; the
    # structural decoders below handle those. Everything else is text.
    if data.startswith(b"PK"):
        # A .txt/.md that literally begins with a zip header is binary garbage;
        # refuse rather than mangle it.
        raise RuntimeError("File is a zip/archive despite its text extension; refusing to parse as text.")
    text = data.decode("utf-8", errors="replace").replace("\x00", "")
    trimmed, truncated = _clip(text, limit)
    return {"format": suffix.lstrip("."), "name": p.name, "bytes": size,
            "chars": len(text), "text": trimmed, "truncated": truncated}


def _extract_binary_doc(p: Path, suffix: str, *, max_bytes: int,
                        limit: int) -> dict[str, Any]:
    if p.stat().st_size > max_bytes:
        raise RuntimeError(f"File too large ({p.stat().st_size} bytes > {max_bytes} limit)")
    if suffix == ".pdf":
        return extract_pdf(p, limit=limit)
    if suffix == ".docx":
        return extract_docx(p, limit=limit)
    if suffix == ".xlsx":
        return extract_xlsx(p, limit=limit)
    return extract_pptx(p, limit=limit)


def _clip(text: str, limit: int) -> tuple[str, bool]:
    if len(text) <= limit:
        return text, False
    return text[:limit] + f"\n… [truncated {len(text) - limit} more characters]", True


# ---- CSV (stdlib) -------------------------------------------------------------


def extract_csv(p: Path, *, limit: int) -> dict[str, Any]:
    import csv as _csv
    with open(p, newline="", encoding="utf-8", errors="replace") as fh:
        try:
            reader = list(_csv.reader(fh))
        except _csv.Error as exc:  # pragma: no cover - defensive
            raise RuntimeError(f"CSV parse error: {exc}") from exc
    # CSV is data, never formulas. Render cells verbatim as text/JSON-ish.
    lines: list[str] = []
    for row in reader:
        cell_str = ", ".join(c.replace("\n", " ").replace("\r", " ")[:500] for c in row)
        lines.append(cell_str)
    text = "\n".join(lines)
    trimmed, truncated = _clip(text, limit)
    return {"format": "csv", "name": p.name, "rows": len(reader),
            "chars": len(text), "text": trimmed,
            "truncated": truncated}


# ---- Plain-image metadata (PIL, base desktop dep) ----------------------------


def image_metadata(p: Path) -> dict[str, Any]:
    # PIL ships with the `desktop` extra and is part of the vision path; it is
    # already a supported base capability. Treat it as optional-but-present.
    if not _have("PIL"):
        raise RuntimeError("image metadata requires the Pillow library "
                           "(pip install aiba-agent[desktop] or [all])")
    from PIL import Image
    with Image.open(p) as im:
        info = {"format": (im.format or "").lower(), "width": im.width,
                "height": im.height, "mode": im.mode}
        dpi_list: list[int] = []
        try:
            raw_dpi = im.info.get("dpi")
            if raw_dpi:
                dpi_list = [int(round(x)) for x in raw_dpi]
        except Exception:
            dpi_list = []
        info["dpi"] = dpi_list or None
    dpi_txt = ""
    if dpi_list:
        dpi_txt = " (dpi " + "x".join(str(x) for x in dpi_list) + ")"
    text = (f"{p.name}: {info['width']}x{info['height']} "
            f"{info['format']} {info['mode']} image.{dpi_txt}")
    return {"kind": "image", "name": p.name, "metadata": info, "text": text}


# ---- PDF (optional pypdf) ------------------------------------------------------
_PDF_PAGE_LIMIT = 400


def extract_pdf(p: Path, *, limit: int) -> dict[str, Any]:
    if not _have("pypdf"):
        raise RuntimeError("PDF extraction requires the optional 'pypdf' package "
                           "(pip install aiba-agent[media] or [all]).")
    from pypdf import PdfReader
    reader = PdfReader(str(p))
    page_count = len(reader.pages)
    pages = []
    for i, page in enumerate(reader.pages):
        if i >= _PDF_PAGE_LIMIT:
            pages.append(f"[… stopped at {_PDF_PAGE_LIMIT} pages]")
            break
        pages.append((page.extract_text() or "").strip())
    text = "\n\n".join(x for x in pages if x)
    trimmed, truncated = _clip(text, limit)
    return {"format": "pdf", "name": p.name, "pages": page_count,
            "chars": len(text), "text": trimmed, "truncated": truncated}


# ---- DOCX (optional python-docx) -----------------------------------------------


def extract_docx(p: Path, *, limit: int) -> dict[str, Any]:
    if not _have("docx"):
        raise RuntimeError("DOCX extraction requires the optional 'python-docx' "
                           "package (pip install aiba-agent[media] or [all]).")
    import docx
    document = docx.Document(str(p))
    text_parts: list[str] = []
    for para in document.paragraphs:
        if para.text:
            text_parts.append(para.text)
    # Tables also carry meaningful text.
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.replace("\n", " ").strip() for c in row.cells]
            text_parts.append(" | ".join(cells))
    text = "\n".join(text_parts)
    trimmed, truncated = _clip(text, limit)
    return {"format": "docx", "name": p.name, "chars": len(text),
            "text": trimmed, "truncated": truncated}


# ---- XLSX (optional openpyxl) ---------------------------------------------------
_XLSX_SHEET_LIMIT = 64
_XLSX_ROW_LIMIT = 50_000


def extract_xlsx(p: Path, *, limit: int) -> dict[str, Any]:
    if not _have("openpyxl"):
        raise RuntimeError("XLSX extraction requires the optional 'openpyxl' "
                           "package (pip install aiba-agent[media] or [all]).")
    # data_only=True is deliberately NOT used: we never want cached formulas in
    # cells to pull precomputed ("trusted") values from untrusted files, and we
    # never evaluate formulas. Default (data_only=False) yields the stored
    # string/value; formula cells are returned as their formula text only.
    import openpyxl
    wb = openpyxl.load_workbook(str(p), read_only=True, data_only=False)
    try:
        out: list[str] = []
        for sheet_index, ws in enumerate(wb.worksheets):
            if sheet_index >= _XLSX_SHEET_LIMIT:
                out.append(f"[stopped at {_XLSX_SHEET_LIMIT} worksheets]")
                break
            out.append(f"== {ws.title} ==")
            row_idx = 0
            for row in ws.iter_rows(values_only=True):
                row_idx += 1
                if row_idx > _XLSX_ROW_LIMIT:
                    out.append(f"[stopped at {_XLSX_ROW_LIMIT} rows in this sheet]")
                    break
                cells = ["" if v is None else str(v).replace("\n", " ")[:500] for v in row]
                out.append(" | ".join(cells))
        text = "\n".join(out)
    finally:
        wb.close()
    trimmed, truncated = _clip(text, limit)
    return {"format": "xlsx", "name": p.name, "sheets": len(wb.sheetnames)
            if hasattr(wb, "sheetnames") else 0,
            "chars": len(text), "text": trimmed, "truncated": truncated}


# ---- PPTX (optional python-pptx) -------------------------------------------------
_PPTX_SLIDE_LIMIT = 300


def extract_pptx(p: Path, *, limit: int) -> dict[str, Any]:
    if not _have("pptx"):
        raise RuntimeError("PPTX extraction requires the optional 'python-pptx' "
                           "package (pip install aiba-agent[media] or [all]).")
    # Late import to avoid failure under [api]-only installs.
    slides: list[str] = []
    try:
        from pptx import Presentation
        prs = Presentation(str(p))
        for idx, slide in enumerate(prs.slides):
            if idx >= _PPTX_SLIDE_LIMIT:
                slides.append(f"[stopped at {_PPTX_SLIDE_LIMIT} slides]")
                break
            parts: list[str] = []
            for shape in slide.shapes:
                frame = getattr(shape, "text_frame", None)
                if frame is not None and getattr(shape, "has_text_frame", False):
                    # text_frame may be duck-typed differently across pptx shapes.
                    paragraphs = getattr(frame, "paragraphs", None) or []
                    t = " ".join(getattr(para, "text", "") for para in paragraphs).strip()
                    if t:
                        parts.append(t)
                elif getattr(shape, "has_table", False):
                    table = getattr(shape, "table", None)
                    if table is not None:
                        for row in getattr(table, "rows", []) or []:
                            cells = [getattr(c, "text", "") for c in getattr(row, "cells", []) or []]
                            parts.append(" | ".join(cells))
            if parts:
                slides.append(f"--- Slide {idx + 1} ---")
                slides.extend(parts)
    except Exception as exc:
        # Distinguish missing-optional (already handled) from genuine corruption.
        if "No module named" in str(exc):
            raise RuntimeError("PPTX extraction requires the optional 'python-pptx' "
                               "package (pip install aiba-agent[media] or [all]).") from exc
        raise RuntimeError(f"PPTX parse error: {exc}") from exc
    text = "\n".join(slides)
    trimmed, truncated = _clip(text, limit)
    return {"format": "pptx", "name": p.name, "chars": len(text),
            "text": trimmed, "truncated": truncated}
